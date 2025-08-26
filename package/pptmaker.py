import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches
import json
import re
API_KEY = st.secrets["GEMINI_API_KEY"]
genai.configure(api_key=API_KEY)
model = genai.GenerativeModel("gemini-2.5-flash")  # safest for now

def generate_ppt_content(topic: str):
    prompt = f"""
    Generate presentation slides in **valid JSON only** for topic: {topic}.
    Strict JSON format:
    {{
      "slides": [
        {{"title": "Slide Title", "points": ["point 1", "point 2"]}},
        {{"title": "Another Title", "points": ["point 1", "point 2"]}}
      ]
    }}
    """

    response = model.generate_content(prompt)
    text = response.text.strip()

    # Try to extract valid JSON
    try:
        data = json.loads(text)
    except:
        match = re.search(r"\{.*\}", text, re.S)
        if match:
            data = json.loads(match.group())
        else:
            st.error("⚠️ Gemini did not return valid JSON.")
            return None

    return data

# --- Function to build PPT file ---
def build_ppt(data, filename="output.pptx"):
    prs = Presentation()
    for slide_data in data.get("slides", []):
        slide_layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slide_data.get("title", "Untitled Slide")
        content.text = "\n".join(slide_data.get("points", []))

    prs.save(filename)
    return filename

def call():
    st.title("📑 Gemini PPT Generator")

    topic = st.text_input("Enter topic for slides:", placeholder="e.g. Artificial Intelligence")

    if st.button("Generate PPT"):
        if topic.strip():
            with st.spinner("Generating slides..."):
                data = generate_ppt_content(topic)

                if data:
                    filename = build_ppt(data)
                    with open(filename, "rb") as f:
                        st.download_button(
                            label="⬇️ Download PPT",
                            data=f,
                            file_name=filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
        else:
            st.warning("Please enter a topic!")
